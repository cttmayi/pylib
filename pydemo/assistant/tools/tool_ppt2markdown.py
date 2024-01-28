from pptx import Presentation
import markdown

def ppt_to_markdown(ppt_file, markdown_file):
    prs = Presentation(ppt_file)
    slides = prs.slides

    markdown_content = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        markdown_content += run.text
                    markdown_content += "\n"

    with open(markdown_file, "w", encoding="utf-8") as file:
        file.write(markdown.markdown(markdown_content))

# 使用示例
ppt_to_markdown("input.pptx", "output.md")
